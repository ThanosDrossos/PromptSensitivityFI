"""Build the 2026-08-03 supervisor deck (graph-heavy, minimal text).

Nine 16:9 slides: recap -> design -> three results -> deliverable -> what it
contributes -> status. Every figure comes from make_supervisor_figures.py, so
the deck cannot drift from the data.

    uv run python -m prompt_sensitivity.scripts.make_supervisor_deck
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt
from loguru import logger

from ..config import load_config
from ..logging_setup import configure_logging

W, H = Inches(13.333), Inches(7.5)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0x21, 0x66, 0xAC)
GREEN = RGBColor(0x1B, 0x78, 0x37)
FIGDIR = "figures/supervisor_2026-08-03"


def _txt(slide, x, y, w, h, text, size, *, bold=False, color=INK, align=1, italic=False):
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT][align]
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            run.font.name = "Segoe UI"
    return tb


def _slide(prs, title, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])          # blank
    _txt(s, Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.7),
         title, 30, bold=True, align=0)
    if subtitle:
        _txt(s, Inches(0.58), Inches(0.98), Inches(12.2), Inches(0.45),
             subtitle, 15.5, color=MUTED, align=0)
    return s


def _pic(slide, root, name, top, height=None, width=None):
    """Insert a figure, horizontally centred."""
    path = str(root / FIGDIR / f"{name}.png")
    if height is not None:
        pic = slide.shapes.add_picture(path, Inches(0), top, height=height)
    else:
        pic = slide.shapes.add_picture(path, Inches(0), top, width=width)
    pic.left = Emu(int((W - pic.width) / 2))
    return pic


def _takeaway(slide, text, top=Inches(6.62), color=ACCENT):
    bar = slide.shapes.add_textbox(Inches(0.55), top, Inches(12.2), Inches(0.6))
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        run.font.size = Pt(17)
        run.font.bold = True
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
    return bar


def build(root, out_path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ---- 1 title -----------------------------------------------------------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _txt(s, Inches(0.9), Inches(2.25), Inches(11.5), Inches(1.2),
         "How much does phrasing decide\nwhether an LLM answers correctly?", 40, bold=True, align=0)
    _txt(s, Inches(0.95), Inches(4.15), Inches(11.5), Inches(0.9),
         "Final results  ·  Thanos Drossos  ·  3 August 2026", 20, color=MUTED, align=0)
    _txt(s, Inches(0.95), Inches(4.85), Inches(11.5), Inches(0.9),
         "Measuring prompt sensitivity in bits — 3 open models, 149 ambiguous questions,\n"
         "all data collection complete", 16, color=MUTED, align=0)

    # ---- 2 recap -----------------------------------------------------------
    s = _slide(prs, "Recap: the question we set out to answer",
               "Users phrase the same question differently — and models reward some phrasings over others.")
    rows = [
        ("The problem",
         "The same question, asked two ways, can flip an LLM from right to wrong.\n"
         "Nobody can say how much of that is the model's knowledge vs. the wording."),
        ("The idea",
         "Borrow Functional Information from biology (Szostak/Hazen):\n"
         "bits = −log₂(fraction of possibilities that still work)."),
        ("The experiment",
         "Take real ambiguous questions, make them measurably more specific,\n"
         "and watch what happens — on three open models."),
    ]
    y = Inches(1.85)
    for head, body in rows:
        _txt(s, Inches(0.75), y, Inches(3.0), Inches(0.9), head, 20, bold=True, color=ACCENT, align=0)
        _txt(s, Inches(4.0), y, Inches(8.6), Inches(1.1), body, 17, align=0)
        y += Inches(1.55)
    _takeaway(s, "Status: everything measured. Today = the results.")

    # ---- 3 design ----------------------------------------------------------
    s = _slide(prs, "What we did", "AmbigQA: humans already wrote both versions of every question — we never invent data.")
    _pic(s, root, "design", Inches(1.65), height=Inches(4.6))
    _takeaway(s, "One dial we control (question specificity) → three things we measure.")

    # ---- 4 headline --------------------------------------------------------
    s = _slide(prs, "Result 1 — specificity buys robustness",
               "Same questions, same evidence, only the wording gets more specific (+1.6 bits).")
    _pic(s, root, "headline", Inches(1.6), height=Inches(4.55))
    _takeaway(s, "Accuracy ≈ doubles, phrasing-luck drops ~0.8 bits, answers converge — in all three models.")

    # ---- 5 evidence dial ---------------------------------------------------
    s = _slide(prs, "Result 2 — but only when the model can find the answer",
               "Same experiment repeated with no / half / full supporting evidence.")
    _pic(s, root, "dial", Inches(1.6), height=Inches(4.5))
    _takeaway(s, "No evidence → being specific buys nothing. The two levers multiply, they don't add.")

    # ---- 6 independence ----------------------------------------------------
    s = _slide(prs, "Result 3 — the three measures are genuinely different",
               "Correlation between every pair of metrics we computed (darker = more redundant).")
    _pic(s, root, "independence", Inches(1.42), height=Inches(5.05))
    _takeaway(s, "Three tight blocks, near-zero between them → you cannot infer one axis from another.")

    # ---- 7 posix -----------------------------------------------------------
    s = _slide(prs, "Result 4 — what the literature actually measures",
               "POSIX (Chatterjee et al. 2024) is published as a “prompt sensitivity index”.")
    _pic(s, root, "posix", Inches(1.6), height=Inches(4.5))
    _takeaway(s, "It lands in the answer-scatter family in all 3 models — our phrasing axis stays unoccupied.")

    # ---- 8 deliverable -----------------------------------------------------
    s = _slide(prs, "The deliverable — a prompt checker",
               "A small linear probe on the model's own internal state: one forward pass, before it answers.")
    _pic(s, root, "feedback", Inches(1.6), height=Inches(4.5))
    _takeaway(s, "Works on questions it never saw, labelled by humans, not by our procedure (AUROC .66 vs .46 baseline).")

    # ---- 9 contributions ---------------------------------------------------
    s = _slide(prs, "What this paper could contribute",
               "Four claims we can now defend with data.")
    items = [
        ("1.  A new axis, not a new score",
         "Phrasing sensitivity (ρ_F) is measurable and provably separate from ability and\n"
         "answer scatter — the first prompt metric that isn't accuracy or entropy in disguise."),
        ("2.  Functional Information transfers out of biology",
         "First application of Szostak/Hazen FI to prompts: sensitivity in bits, on a scale\n"
         "that is comparable across questions and models."),
        ("3.  A tidying result for the field",
         "Several published “sensitivity” indices measure one construct — answer dispersion.\n"
         "Our correlation map shows which metric belongs where."),
        ("4.  A usable artefact",
         "“Your prompt is too vague” — predicted before generating, ~10× cheaper than sampling,\n"
         "and it survives a human-labelled held-out test."),
    ]
    y = Inches(1.75)
    for head, body in items:
        _txt(s, Inches(0.75), y, Inches(11.9), Inches(0.4), head, 19, bold=True, color=GREEN, align=0)
        _txt(s, Inches(1.05), y + Inches(0.42), Inches(11.6), Inches(0.75), body, 15.5, align=0)
        y += Inches(1.28)
    _takeaway(s, "Open question for today: which of these should lead the paper?", top=Inches(6.75))

    # ---- 10 status ---------------------------------------------------------
    s = _slide(prs, "Where we stand", None)
    done = [
        "3 models × 149 questions × 2 specificity levels",
        "Evidence dial, k=20 stability, POSIX — all complete",
        "Independence: correlations, factors, counterexamples",
        "Prompt-checker probes + held-out human test",
    ]
    todo = [
        "Write the paper (all numbers are in place)",
        "Decide the lead contribution",
        "Optional: second dataset for external validity",
    ]
    _txt(s, Inches(0.75), Inches(1.5), Inches(6.0), Inches(0.5), "Done — data collection is closed",
         21, bold=True, color=GREEN, align=0)
    y = Inches(2.15)
    for t in done:
        _txt(s, Inches(0.95), y, Inches(5.9), Inches(0.75), "•  " + t, 15.5, align=0)
        y += Inches(0.82)
    _txt(s, Inches(7.3), Inches(1.5), Inches(5.4), Inches(0.5), "Next", 21, bold=True, color=ACCENT, align=0)
    y = Inches(2.15)
    for t in todo:
        _txt(s, Inches(7.5), y, Inches(5.2), Inches(0.75), "•  " + t, 15.5, align=0)
        y += Inches(0.82)
    _takeaway(s, "Ask: feedback on framing + the lead contribution, before drafting.")

    prs.save(str(out_path))
    return len(prs.slides._sldIdLst)


def main() -> int:
    configure_logging("make_supervisor_deck")
    root = load_config().repo_root()
    out = root.parent.parent / "Prompt_Sensitivity_Final_Results_2026-08-03.pptx"
    n = build(root, out)
    logger.info("wrote {} ({} slides)", out.name, n)
    print(f"DONE {out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
