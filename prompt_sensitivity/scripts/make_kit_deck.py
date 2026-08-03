"""Rebuild the final-results section of the deck INSIDE the KIT template.

The 2026-08-03 results slides were pasted into
`Prompt_Sensitivity_Full_Results_final.pptx` on the "Blank" layout, i.e.
without KIT chrome. This script drops those slides and re-creates them on the
KIT master ("Title and Text" layout -> KIT title styling, logo, footer, slide
number), adds the two technical slides (metric definitions with formulas, and
the probe architecture), and keeps every figure sourced from the generators so
the deck cannot drift from the data.

A timestamped backup of the original file is written first.

    uv run python -m prompt_sensitivity.scripts.make_kit_deck
"""

from __future__ import annotations

import argparse
import shutil

from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ..config import load_config
from ..logging_setup import configure_logging

DECK = "Prompt_Sensitivity_Full_Results_final.pptx"
FIRST_NEW_SLIDE = 23           # 1-based: slides 23..end are the pasted ones
LAYOUT_TITLE_TEXT = 13         # "Title and Text" — carries the KIT header/footer

# KIT corporate colours (Gestaltungsrichtlinien).
KIT_GREEN = RGBColor(0x00, 0x96, 0x82)
KIT_BLUE = RGBColor(0x46, 0x64, 0xAA)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
FONT = "Arial"                 # KIT corporate typeface

W = Inches(13.333)
FIG_TOP = Inches(1.62)         # below the KIT title band
TAKE_TOP = Inches(6.16)        # above the KIT footer band (starts 6.86)
FIGDIR = "figures/supervisor_2026-08-03"


# --------------------------------------------------------------------------- #
# helpers                                                                     #
# --------------------------------------------------------------------------- #


def drop_slides_from(prs, first_1based: int) -> int:
    """Remove slides [first_1based .. end] and their relationships."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    removed = 0
    for sld in slides[first_1based - 1:]:
        rId = sld.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.part.drop_rel(rId)
        xml_slides.remove(sld)
        removed += 1
    return removed


def _style(run, size, bold=False, color=INK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT


def kit_slide(prs, title: str, subtitle: str | None = None):
    """New slide on the KIT 'Title and Text' layout; body placeholder removed
    (we position content ourselves), title placeholder keeps KIT styling."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_TEXT])
    slide.shapes.title.text_frame.text = title
    for p in slide.shapes.title.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = FONT
    # Drop the empty body placeholder so it cannot render "click to add text".
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == 13:
            ph._element.getparent().remove(ph._element)
    if subtitle:
        tb = slide.shapes.add_textbox(Inches(0.41), Inches(1.20), Inches(12.4), Inches(0.4))
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.text = subtitle
        for r in p.runs:
            _style(r, 14, color=MUTED)
    return slide


def add_figure(slide, root, name, top=FIG_TOP, height=Inches(4.3)):
    pic = slide.shapes.add_picture(str(root / FIGDIR / f"{name}.png"),
                                   Inches(0), top, height=height)
    pic.left = Emu(int((W - pic.width) / 2))
    return pic


def add_takeaway(slide, text, top=TAKE_TOP, color=KIT_GREEN):
    tb = slide.shapes.add_textbox(Inches(0.41), top, Inches(12.4), Inches(0.55))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    for r in p.runs:
        _style(r, 15, bold=True, color=color)
    return tb


def add_body(slide, blocks, top=Inches(1.75), left=Inches(0.41),
             width=Inches(12.4), gap=Inches(1.16), head_size=17, body_size=13.5):
    """[(heading, text)] rendered as KIT-blue heading + dark body."""
    y = top
    for head, text in blocks:
        tb = slide.shapes.add_textbox(left, y, width, Inches(0.36))
        p = tb.text_frame.paragraphs[0]
        p.text = head
        for r in p.runs:
            _style(r, head_size, bold=True, color=KIT_BLUE)
        tb2 = slide.shapes.add_textbox(left + Inches(0.22), y + Inches(0.38),
                                       width - Inches(0.22), Inches(0.7))
        tb2.text_frame.word_wrap = True
        for i, line in enumerate(text.split("\n")):
            p2 = tb2.text_frame.paragraphs[0] if i == 0 else tb2.text_frame.add_paragraph()
            p2.text = line
            p2.alignment = PP_ALIGN.LEFT
            for r in p2.runs:
                _style(r, body_size)
        y += gap
    return y


# --------------------------------------------------------------------------- #
# the results section                                                         #
# --------------------------------------------------------------------------- #


def build_section(prs, root):
    # 1 · recap ------------------------------------------------------------
    s = kit_slide(prs, "Recap: the question we set out to answer",
                  "Users phrase the same question differently — and models reward some phrasings over others.")
    add_body(s, [
        ("The problem",
         "The same question, asked two ways, can flip an LLM from right to wrong.\n"
         "Nobody can say how much of that is the model's knowledge vs. the wording."),
        ("The idea — Functional Information (Szostak 2003, Hazen et al. 2007)",
         "Measure everything in bits:  bits = −log₂(fraction of possibilities that still work).\n"
         "Imported from biology, where it scores how rare a functional molecule is."),
        ("The experiment",
         "Take real ambiguous questions, make them measurably more specific,\n"
         "and watch three independent quantities respond — on three open 7–8B models."),
    ])
    add_takeaway(s, "Data collection is complete. Today: the results.")

    # 2 · design -----------------------------------------------------------
    s = kit_slide(prs, "What we did",
                  "AmbigQA (Min et al. 2020): human annotators already wrote both versions of every question.")
    add_figure(s, root, "design", height=Inches(4.25))
    add_takeaway(s, "One dial we control (FI_spec, model-free) → three things we measure.")

    # 3 · definitions (NEW, technical) -------------------------------------
    s = kit_slide(prs, "The four quantities, defined",
                  "Same ruler throughout — only the space being counted changes.")
    add_figure(s, root, "definitions", top=Inches(1.55), height=Inches(4.85))
    add_takeaway(s, "AUFI_in summarises the FI_in(k) curve in one number: the bits of rephrasing-rarity a question demands.",
                 top=Inches(6.42))

    # 4 · result 1 ---------------------------------------------------------
    s = kit_slide(prs, "Result 1 — effect of question specificity",
                  "Paired within question: identical evidence and gold answer; only the question text changes "
                  "(FI_spec 0 → 1.58 bits). n = 149 × 3 models.")
    add_figure(s, root, "headline", height=Inches(4.3))
    add_takeaway(s, "ΔF̄ = +0.22…+0.25 · ΔAUFI_in = −0.79…−0.87 bits · ΔH_sem < 0 — "
                    "all 3 models, Wilcoxon signed-rank p ≤ 5e-9.")

    # 5 · result 2 ---------------------------------------------------------
    s = kit_slide(prs, "Result 2 — specificity × evidence interaction",
                  "Same 50 questions re-run with 0 %, 50 % and 100 % of the retrieved evidence snippets "
                  "(identical across both specificity levels).")
    add_figure(s, root, "dial", height=Inches(4.3))
    add_takeaway(s, "ΔF̄ = +0.02 / +0.10 / +0.24 and ΔAUFI_in = −0.05 / −0.31 / −0.83 bits "
                    "at 0 / 50 / 100 % evidence — the two factors interact.")

    # 6 · result 3 ---------------------------------------------------------
    s = kit_slide(prs, "Result 3 — the three measures are genuinely different",
                  "Within-level Spearman |ρ| between all 12 metrics we computed (n = 149 questions × 3 models).")
    add_figure(s, root, "independence", top=Inches(1.5), height=Inches(4.85))
    add_takeaway(s, "ρ_F ⊥ accuracy (.08) and ⊥ H_sem (.03) → no axis can be inferred from another.",
                 top=Inches(6.42))

    # 7 · result 4 ---------------------------------------------------------
    s = kit_slide(prs, "Result 4 — what the literature actually measures",
                  "POSIX (Chatterjee et al. 2024): a published prompt-sensitivity index, computed on our data.")
    add_figure(s, root, "posix", height=Inches(4.3))
    add_takeaway(s, "It loads on answer dispersion in all 3 models (ρ .35–.69) — the phrasing axis stays empty.")

    # 8 · probe design (NEW, technical) ------------------------------------
    s = kit_slide(prs, "How the prompt checker works",
                  "A linear probe on the model's own hidden state — no fine-tuning, no extra annotation.")
    add_figure(s, root, "probe", top=Inches(1.5), height=Inches(4.8))
    add_takeaway(s, "Cost: one forward pass — the warning arrives BEFORE the model answers.",
                 top=Inches(6.42))

    # 9 · probe results ----------------------------------------------------
    s = kit_slide(prs, "Prompt checker — does it work?",
                  "Out-of-fold by question; right panel = 1,852 questions never seen, labelled by AmbigQA annotators.")
    add_figure(s, root, "feedback", height=Inches(4.3))
    add_takeaway(s, "Vagueness .85–.87 in-distribution, .66 on unseen human-labelled questions (length baseline .46).")

    # 10 · contributions ---------------------------------------------------
    s = kit_slide(prs, "What this paper could contribute",
                  "Four claims we can now defend with data.")
    add_body(s, [
        ("1 · A new axis, not another score",
         "ρ_F is measurably separate from ability and dispersion — the first prompt-sensitivity\n"
         "metric that is not accuracy or entropy in disguise (validated by ICC + factor analysis)."),
        ("2 · Functional Information transfers out of biology",
         "First application of the Szostak/Hazen formalism to prompts: sensitivity in bits,\n"
         "comparable across questions and models."),
        ("3 · A tidying result for the field",
         "POSIX, S_τ, TVD, |A_q| and variation ratio all load on ONE construct — answer dispersion.\n"
         "The correlation map shows which published metric belongs where."),
        ("4 · A usable artefact",
         "“Your prompt is too vague”, predicted from one forward pass, ~10× cheaper than sampling,\n"
         "and it survives a human-labelled held-out test."),
    ], top=Inches(1.72), gap=Inches(1.15), head_size=16, body_size=13)
    add_takeaway(s, "Open question for today: which of these should lead the paper?", top=Inches(6.5))

    # 11 · status ----------------------------------------------------------
    s = kit_slide(prs, "Where we stand", None)
    left = [
        "3 models × 149 questions × 2 specificity levels",
        "Evidence dial, k=20 stability, POSIX — all complete",
        "Independence: correlations, factors, counterexamples",
        "Prompt-checker probes + held-out human test",
    ]
    right = [
        "Write the paper (all numbers are in place)",
        "Decide the lead contribution",
        "Optional: second dataset for external validity",
    ]
    for x, head, items, col in [(Inches(0.41), "Done — data collection is closed", left, KIT_GREEN),
                               (Inches(7.0), "Next", right, KIT_BLUE)]:
        tb = s.shapes.add_textbox(x, Inches(1.75), Inches(6.2), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = head
        for r in p.runs:
            _style(r, 18, bold=True, color=col)
        y = Inches(2.42)
        for it in items:
            b = s.shapes.add_textbox(x + Inches(0.18), y, Inches(6.0), Inches(0.5))
            b.text_frame.word_wrap = True
            p = b.text_frame.paragraphs[0]
            p.text = "•  " + it
            for r in p.runs:
                _style(r, 13.5)
            y += Inches(0.72)
    add_takeaway(s, "Ask: feedback on the framing and on the lead contribution, before drafting.")


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--deck", type=str, default=DECK)
    ap.add_argument("--no-backup", action="store_true")
    args = ap.parse_args()

    configure_logging("make_kit_deck")
    root = load_config().repo_root()
    seminar = root.parent.parent
    path = seminar / args.deck
    if not path.exists():
        logger.error("deck not found: {}", path)
        return 1
    if not args.no_backup:
        bak = path.with_name(path.stem + "_backup.pptx")
        shutil.copy2(path, bak)
        logger.info("backup -> {}", bak.name)

    prs = Presentation(str(path))
    n_before = len(prs.slides._sldIdLst)
    removed = drop_slides_from(prs, FIRST_NEW_SLIDE)
    logger.info("removed {} pasted slides (kept {})", removed, FIRST_NEW_SLIDE - 1)
    build_section(prs, root)
    prs.save(str(path))
    n_after = len(prs.slides._sldIdLst)
    logger.info("{}: {} -> {} slides", path.name, n_before, n_after)
    print(f"DONE {path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
